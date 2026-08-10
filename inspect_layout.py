import pptx

prs = pptx.Presentation('3C時代不抓狂的親子相處之道_Google簡報相容版 202608010.pptx')
slide_w = prs.slide_width
slide_h = prs.slide_height

with open('pptx_exact_layout_utf8.txt', 'w', encoding='utf-8') as f:
    f.write(f"Slide WxH: {slide_w} x {slide_h}\n")
    for idx, slide in enumerate(prs.slides, 1):
        f.write(f"\n==================== SLIDE {idx} ====================\n")
        for s_idx, shape in enumerate(slide.shapes, 1):
            l = round(shape.left / slide_w * 100, 1) if shape.left else 0
            t = round(shape.top / slide_h * 100, 1) if shape.top else 0
            w = round(shape.width / slide_w * 100, 1) if shape.width else 0
            h = round(shape.height / slide_h * 100, 1) if shape.height else 0
            stype = str(shape.shape_type)
            f.write(f"Shape {s_idx} [{stype}] name='{shape.name}' -> L:{l}% T:{t}% W:{w}% H:{h}%\n")
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = p.text.strip()
                    if txt:
                        font_size = p.font.size.pt if p.font and p.font.size else "auto"
                        font_bold = p.font.bold if p.font else None
                        color = "default"
                        try:
                            if p.font and p.font.color and p.font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB:
                                color = str(p.font.color.rgb)
                        except: pass
                        f.write(f"   [Text] '{txt}' (size={font_size}, bold={font_bold}, color={color})\n")
